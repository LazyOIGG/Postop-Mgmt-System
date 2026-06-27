"""
全周期健康管理系统汇报PPT生成脚本
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 配色方案 ──────────────────────────────────────────
PRIMARY    = RGBColor(0x1A, 0x56, 0xDB)  # 主蓝
SECONDARY  = RGBColor(0x22, 0x71, 0xF1)  # 浅蓝
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)  # 青蓝
DARK       = RGBColor(0x1E, 0x29, 0x3B)  # 深灰
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── 工具函数 ──────────────────────────────────────────
def add_blank():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, width, height, color, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def txt(slide, left, top, width, height, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='微软雅黑'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def multi_txt(slide, left, top, width, height, lines, size=14, color=DARK, line_space=Pt(28)):
    """lines: list of str or (str, bool) tuples for bold"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = '微软雅黑'
        p.space_after = line_space
    return txBox

def title_bar(slide, title, subtitle=""):
    """通用标题栏"""
    rect(slide, 0, 0, SW, Inches(1.1), PRIMARY)
    txt(slide, Inches(0.8), Inches(0.15), Inches(8), Inches(0.5),
        title, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, Inches(0.8), Inches(0.65), Inches(8), Inches(0.35),
            subtitle, size=13, color=RGBColor(0xBF, 0xDB, 0xFE))

def section_header(slide, title, icon=""):
    """页面内小节标题"""
    txt(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
        f"{icon}  {title}" if icon else title,
        size=22, bold=True, color=PRIMARY)

def card(slide, left, top, width, height, title, items, icon="", title_color=PRIMARY):
    """内容卡片"""
    # 卡片底
    shp = rect(slide, left, top, width, height, WHITE)
    # 顶部色条
    rect(slide, left, top, width, Inches(0.05), title_color)
    # 标题
    txt(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.35),
        f"{icon}  {title}" if icon else title, size=15, bold=True, color=title_color)
    # 内容
    multi_txt(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.5), height - Inches(0.6),
              items, size=12, color=GRAY, line_space=Pt(22))


# ══════════════════════════════════════════════════════════
#  Slide 1 — 封面
# ══════════════════════════════════════════════════════════
s = add_blank()
rect(s, 0, 0, SW, SH, PRIMARY)
# 装饰圆
circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
circle.fill.solid(); circle.fill.fore_color.rgb = SECONDARY; circle.line.fill.background()
circle2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(4.5), Inches(4), Inches(4))
circle2.fill.solid(); circle2.fill.fore_color.rgb = RGBColor(0x15, 0x4A, 0xC9); circle2.line.fill.background()

txt(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1.0),
    "术后全周期健康管理系统", size=44, bold=True, color=WHITE)
txt(s, Inches(1.5), Inches(3.2), Inches(10), Inches(0.6),
    "基于多智能体与知识图谱的智能术后康复管理平台", size=20, color=RGBColor(0xBF, 0xDB, 0xFE))
txt(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.4),
    "Postop Full-Cycle Health Management System", size=14, color=RGBColor(0x93, 0xBB, 0xFD))

# 底部信息
txt(s, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4),
    "软件工程创新实践", size=16, color=RGBColor(0xBF, 0xDB, 0xFE))
txt(s, Inches(1.5), Inches(6.3), Inches(10), Inches(0.4),
    "团队成员：倪春阳  葛兴  桂俊杰  林学成", size=14, color=RGBColor(0x93, 0xBB, 0xFD))


# ══════════════════════════════════════════════════════════
#  Slide 2 — 目录
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "目  录", "CONTENTS")

toc = [
    ("01", "项目介绍",      "项目背景、目标与技术选型"),
    ("02", "系统架构总览",  "前后端架构、技术栈与部署方案"),
    ("03", "核心功能演示",  "多智能体对话、康复计划、知识图谱等"),
    ("04", "个人分工",      "四位成员的职责与贡献"),
    ("05", "项目总结",      "成果、亮点与未来规划"),
]
for i, (num, title, desc) in enumerate(toc):
    y = Inches(1.6) + Inches(i * 1.05)
    # 编号圆
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.6), Inches(0.6))
    circ.fill.solid(); circ.fill.fore_color.rgb = PRIMARY; circ.line.fill.background()
    tf = circ.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.text = num; p.font.size = Pt(18); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER; tf.paragraphs[0].font.name = 'Arial'
    # 标题
    txt(s, Inches(2.1), y - Inches(0.02), Inches(4), Inches(0.4),
        title, size=20, bold=True, color=DARK)
    txt(s, Inches(2.1), y + Inches(0.38), Inches(6), Inches(0.3),
        desc, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 3 — 项目背景
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "项目介绍", "PART 01 · 项目背景")

problems = [
    ("🏥", "术后管理痛点", "传统术后管理依赖患者自述和定期复查，缺乏连续性监测，信息碎片化严重"),
    ("🧠", "智能化需求", "患者需要个性化康复指导、用药提醒和异常预警，现有系统难以满足"),
    ("📊", "数据孤岛", "患者健康数据分散在不同系统，医生难以全面掌握患者状态"),
    ("💬", "医患沟通低效", "术后患者与医生沟通渠道有限，问题响应不及时"),
]
for i, (icon, title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.1)
    y = Inches(1.6) + Inches(row * 2.6)
    card(s, x, y, Inches(5.7), Inches(2.1), title,
         [desc], icon=icon, title_color=SECONDARY)


# ══════════════════════════════════════════════════════════
#  Slide 4 — 项目目标
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "项目介绍", "PART 01 · 项目目标")

goals = [
    ("🤖 智能对话", "基于多智能体架构的医学问答系统，\n支持NER实体识别与知识图谱增强"),
    ("📋 康复管理", "AI驱动的分阶段康复计划生成，\n量化指标追踪与成就激励系统"),
    ("⚠️ 异常预警", "规则引擎+时序异常检测，\n高风险患者自动预警通知"),
    ("🔗 知识图谱", "Neo4j医学知识图谱，\nText2Cypher自然语言查询"),
    ("📱 多端适配", "患者端+医生端双角色架构，\n响应式Web界面"),
    ("🔒 安全可靠", "JWT认证+Token黑名单，\nCypher注入防护"),
]
for i, (title, desc) in enumerate(goals):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(col * 4.1)
    y = Inches(1.5) + Inches(row * 2.7)
    shp = rect(s, x, y, Inches(3.8), Inches(2.3), LIGHT_BLUE)
    shp.line.fill.background()
    txt(s, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.4),
        title, size=16, bold=True, color=PRIMARY)
    txt(s, x + Inches(0.2), y + Inches(0.65), Inches(3.4), Inches(1.5),
        desc, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 5 — 技术选型
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "项目介绍", "PART 01 · 技术选型")

# 后端
card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "后端技术栈", [
         "▸ Web框架: FastAPI + Uvicorn",
         "▸ 数据库: MySQL 8.0 (连接池)",
         "▸ 图数据库: Neo4j 5",
         "▸ 缓存: Redis 7",
         "▸ LLM: DeepSeek API (OpenAI SDK)",
         "▸ NER: BERT + RNN + LoRA微调",
         "▸ 语音: 阿里云DashScope ASR",
         "▸ OCR: PaddleOCR + PaddlePaddle",
         "▸ 安全: JWT + Token黑名单",
         "▸ 容器化: Docker Compose",
     ], icon="⚙️", title_color=PRIMARY)

# 前端
card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "前端技术栈", [
         "▸ 框架: Vue 3 + TypeScript",
         "▸ 构建: Vite 8",
         "▸ UI库: Element Plus (按需引入)",
         "▸ 状态管理: Pinia 3",
         "▸ 路由: Vue Router 4 (懒加载)",
         "▸ HTTP: Axios + 拦截器",
         "▸ 图表: ECharts + vue-echarts",
         "▸ 通信: SSE + WebSocket",
         "▸ 响应式: 媒体查询 + Flex布局",
     ], icon="🎨", title_color=SECONDARY)


# ══════════════════════════════════════════════════════════
#  Slide 6 — 系统架构总览
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "系统架构总览", "PART 02 · 整体架构")

# 架构层级图
layers = [
    ("前端层 (Vue 3)", "患者端: 首页 · AI对话 · 健康打卡 · 康复管理 · 消息\n医生端: 仪表盘 · 患者管理 · 告警 · 统计", SECONDARY),
    ("API网关层", "FastAPI RESTful API (21个端点) + WebSocket + SSE\nJWT认证 · CORS · 结构化日志 · Prometheus监控", PRIMARY),
    ("多智能体层", "CoordinatorAgent意图路由 → 5个专业Agent\nMedicalQA · HealthAssessment · Reminder · Psychology · RehabPlan", RGBColor(0x7C, 0x3A, 0xED)),
    ("服务层", "LLM · NER · 意图识别 · 知识图谱 · Text2Cypher\n健康评估 · 异常检测 · 康复计划 · 多模态(语音/OCR)", DARK),
    ("数据层", "MySQL 8.0 (关系数据) + Neo4j 5 (知识图谱) + Redis 7 (缓存)", GRAY),
]
for i, (title, desc, color) in enumerate(layers):
    y = Inches(1.4) + Inches(i * 1.15)
    shp = rect(s, Inches(0.8), y, Inches(11.7), Inches(1.0), color)
    shp.line.fill.background()
    txt(s, Inches(1.0), y + Inches(0.05), Inches(3), Inches(0.35),
        title, size=15, bold=True, color=WHITE)
    txt(s, Inches(1.0), y + Inches(0.4), Inches(11.3), Inches(0.55),
        desc, size=11, color=RGBColor(0xF0, 0xF0, 0xF0))


# ══════════════════════════════════════════════════════════
#  Slide 7 — 多智能体架构详解
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "系统架构总览", "PART 02 · 多智能体架构")

# 流程图式布局
flow = [
    ("用户输入", Inches(0.3), GRAY),
    ("Coordinator\n意图路由", Inches(2.6), PRIMARY),
    ("", Inches(4.9), None),  # 分叉
]
# 入口
rect(s, Inches(0.3), Inches(2.8), Inches(1.8), Inches(1.2), GRAY)
txt(s, Inches(0.35), Inches(3.0), Inches(1.7), Inches(0.8),
    "用户输入\n文本/语音/图片", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 箭头
arrow1 = rect(s, Inches(2.1), Inches(3.3), Inches(0.5), Inches(0.06), PRIMARY)

# 协调器
rect(s, Inches(2.6), Inches(2.5), Inches(2.2), Inches(1.8), PRIMARY)
txt(s, Inches(2.65), Inches(2.7), Inches(2.1), Inches(1.4),
    "CoordinatorAgent\n意图路由协调器\n(LLM分析)", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 箭头
arrow2 = rect(s, Inches(4.8), Inches(3.3), Inches(0.5), Inches(0.06), PRIMARY)

# 五个Agent
agents = [
    ("MedicalQA\n医学问答", "NER+KG增强"),
    ("Health\n健康评估", "规则引擎+异常检测"),
    ("Reminder\n提醒管理", "工具调用CRUD"),
    ("Psychology\n心理辅导", "LLM对话"),
    ("RehabPlan\n康复计划", "LLM生成+任务"),
]
colors = [RGBColor(0x7C,0x3A,0xED), RGBColor(0x05,0x96,0x69), RGBColor(0xD9,0x77,0x06), RGBColor(0xDC,0x26,0x26), RGBColor(0x25,0x63,0xEB)]
for i, (name, desc) in enumerate(agents):
    x = Inches(5.3) + Inches(i * 1.6)
    rect(s, x, Inches(2.0), Inches(1.4), Inches(1.2), colors[i])
    txt(s, x + Inches(0.05), Inches(2.05), Inches(1.3), Inches(0.6),
        name, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.05), Inches(2.6), Inches(1.3), Inches(0.5),
        desc, size=9, color=RGBColor(0xE0,0xE0,0xE0), align=PP_ALIGN.CENTER)

# 工具系统
rect(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.8), RGBColor(0xFE,0xF3,0xC7))
txt(s, Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.7),
    "🔧 工具系统 (ToolRegistry): 提醒CRUD · 打卡查询 · 药品/疾病查询 · 康复计划 · 天气查询 — 基于OpenAI Function Calling格式",
    size=13, bold=True, color=RGBColor(0x92,0x40,0x0E))

# 数据流
rect(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8), RGBColor(0xEC,0xFD,0xF5))
txt(s, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.7),
    "📊 数据流: MySQL(用户/打卡/消息) + Neo4j(药品/疾病/症状关系) + Redis(会话缓存/Token黑名单)",
    size=13, bold=True, color=RGBColor(0x16,0x65,0x34))


# ══════════════════════════════════════════════════════════
#  Slide 8 — 核心功能：多智能体对话
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "核心功能演示", "PART 03 · 多智能体智能对话")

card(s, Inches(0.6), Inches(1.4), Inches(6.0), Inches(2.5),
     "对话流程", [
         "▸ 用户输入 → CoordinatorAgent意图路由",
         "▸ 专业Agent接收 → 调用LLM + 工具",
         "▸ 支持流式SSE/WebSocket实时响应",
         "▸ 多轮上下文 + Redis缓存 + 摘要压缩",
     ], icon="💬", title_color=PRIMARY)

card(s, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.5),
     "NER实体识别", [
         "▸ BERT + RNN + LoRA微调模型",
         "▸ Aho-Corasick多模式匹配",
         "▸ 实体别名归一化（如：扑热息痛→对乙酰氨基酚）",
         "▸ 支持药物、疾病、症状、食物等7类实体",
     ], icon="🔍", title_color=RGBColor(0x7C,0x3A,0xED))

# 截图占位
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8), LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0),
    "[ 截图区域：AI对话界面演示 ]\n请截取患者端AI对话页面，展示多轮对话效果",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  Slide 9 — 核心功能：康复计划
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "核心功能演示", "PART 03 · 智能康复计划")

features = [
    ("📋 计划生成", "LLM根据手术类型生成\n分阶段康复方案\n急性期→恢复期→强化期"),
    ("📊 指标追踪", "量化康复指标记录\n趋势图表可视化\nECharts动态展示"),
    ("🏋️ 运动指导", "运动库浏览\nExercisePlayer播放\n分阶段运动推荐"),
    ("📝 康复日志", "心情/疼痛/睡眠记录\n文字+图片+语音\n多模态日记"),
    ("🏆 成就系统", "任务完成激励\n成就解锁弹窗\n进度可视化"),
    ("📅 康复日历", "日历视图展示\n任务排期管理\n阶段时间线"),
]
for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(col * 4.1)
    y = Inches(1.4) + Inches(row * 2.8)
    shp = rect(s, x, y, Inches(3.8), Inches(2.4), WHITE)
    rect(s, x, y, Inches(3.8), Inches(0.06), ACCENT)
    txt(s, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.35),
        title, size=16, bold=True, color=PRIMARY)
    txt(s, x + Inches(0.2), y + Inches(0.65), Inches(3.4), Inches(1.6),
        desc, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 10 — 核心功能：知识图谱
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "核心功能演示", "PART 03 · 医学知识图谱")

card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(2.5),
     "图谱结构", [
         "▸ 7种节点: 疾病·药品·食物",
         "  症状·检查项目·科目·治疗方法",
         "▸ 8种关系类型",
         "▸ Neo4j Cypher查询",
     ], icon="🕸️", title_color=PRIMARY)

card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(2.5),
     "Text2Cypher", [
         "▸ LLM动态生成Cypher查询",
         "▸ 自然语言→图查询",
         "▸ 只读安全校验",
         "▸ 禁止CREATE/DELETE操作",
     ], icon="🔗", title_color=RGBColor(0x7C,0x3A,0xED))

card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(2.5),
     "智能搜索网关", [
         "▸ NER实体识别",
         "▸ 意图→KG查询映射",
         "▸ Prompt增强生成",
         "▸ 实体关系可视化",
     ], icon="🔎", title_color=GREEN)

# 截图占位
rect(s, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.8), LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0),
    "[ 截图区域：知识图谱查询与可视化演示 ]\n请截取知识图谱搜索界面，展示实体关系图",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  Slide 11 — 核心功能：健康评估与预警
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "核心功能演示", "PART 03 · 健康评估与异常预警")

card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(2.6),
     "健康评估", [
         "▸ 规则引擎风险分级",
         "  (高/中/低风险关键词匹配)",
         "▸ 体温/血压数值阈值检测",
         "▸ LLM生成结构化健康建议",
         "▸ 多模态输入(文本/语音/图片)",
     ], icon="🏥", title_color=PRIMARY)

card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(2.6),
     "异常检测", [
         "▸ Z-score统计异常检测",
         "▸ 滑动窗口趋势分析",
         "▸ 体温连续3天上升预警",
         "▸ 时序异常自动识别",
         "▸ 高风险患者自动通知",
     ], icon="⚠️", title_color=ORANGE)

card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(2.6),
     "医生端告警", [
         "▸ 实时告警推送",
         "▸ 高风险患者列表",
         "▸ 异常打卡记录",
         "▸ 医患消息通信",
         "▸ 康复方案调整",
     ], icon="👨‍⚕️", title_color=GREEN)

# 截图占位
rect(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.7), LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0),
    "[ 截图区域：医生端仪表盘与告警界面 ]\n请截取医生端Dashboard，展示高风险患者列表和告警信息",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  Slide 12 — 核心功能：多模态交互
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "核心功能演示", "PART 03 · 多模态交互")

card(s, Inches(0.6), Inches(1.4), Inches(3.8), Inches(2.4),
     "语音识别 (ASR)", [
         "▸ 阿里云DashScope Fun-ASR",
         "▸ Whisper模型支持",
         "▸ 患者语音输入→文字",
         "▸ 支持健康打卡语音录入",
     ], icon="🎤", title_color=PRIMARY)

card(s, Inches(4.7), Inches(1.4), Inches(3.8), Inches(2.4),
     "OCR图像识别", [
         "▸ PaddleOCR + PaddlePaddle",
         "▸ 检查报告/处方识别",
         "▸ 图片→结构化数据",
         "▸ 支持健康档案OCR录入",
     ], icon="📷", title_color=SECONDARY)

card(s, Inches(8.8), Inches(1.4), Inches(3.8), Inches(2.4),
     "医患消息通信", [
         "▸ 文本/图片/语音消息",
         "▸ WebSocket实时通信",
         "▸ 通知消息推送",
         "▸ 消息已读状态追踪",
     ], icon="💌", title_color=GREEN)

# 截图占位
rect(s, Inches(0.6), Inches(4.1), Inches(12.1), Inches(2.9), LIGHT_GRAY)
txt(s, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.0),
    "[ 截图区域：多模态交互界面 ]\n请截取语音输入、图片上传、OCR识别等交互界面",
    size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  Slide 13 — 个人分工（倪春阳）
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "个人分工", "PART 04 · 倪春阳")

card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "核心架构与AI智能体", [
         "▸ 术后管理系统重构后初版",
         "▸ 配置规范化、依赖管理",
         "▸ 优化项目整体架构",
         "",
         "▸ 搭建多智能体架构",
         "  新增4个专业Agent并集成到聊天端点",
         "▸ 完成C2任务：模型优化与微调",
         "▸ 实现C2优化方案：实体归一化与多实体识别",
         "",
         "▸ JWT认证重构",
         "  ApiResponse统一响应 + 全局异常处理",
         "▸ 优化注册逻辑，医生患者端面板分离",
     ], icon="🏗️", title_color=PRIMARY)

card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "具体负责模块（请人工补充）", [
         "▸ 后端项目结构搭建与目录规范",
         "",
         "▸ 多智能体编排器 (Orchestrator)",
         "▸ Agent基类与工具调用系统",
         "",
         "▸ 数据库连接池设计",
         "",
         "▸ [请补充其他负责内容]",
         "",
         "",
         "",
         "",
         "",
     ], icon="📝", title_color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 14 — 个人分工（葛兴）
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "个人分工", "PART 04 · 葛兴")

card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "健康管理模块与康复计划", [
         "▸ 健康评估/档案/打卡/趋势分析/提醒中心",
         "▸ 医生端管理（患者列表、高风险记录）",
         "▸ 系统统计与项目展示功能模块",
         "",
         "▸ 康复计划管理模块",
         "  数据库建表、Agent路由、API接口",
         "▸ 康复计划全面升级",
         "  量化指标、运动库、日志、成就系统",
         "▸ 从术后管理扩展为全周期健康管理",
         "",
         "▸ 通知消息模块及医患消息通知",
         "▸ 图片语音交互功能",
     ], icon="🏥", title_color=PRIMARY)

card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "具体负责模块（请人工补充）", [
         "▸ 前端页面开发",
         "",
         "▸ [请补充其他负责内容]",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
     ], icon="📝", title_color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 15 — 个人分工（桂俊杰）
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "个人分工", "PART 04 · 桂俊杰")

card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "知识图谱与安全", [
         "▸ 知识图谱增强与推送通知系统",
         "▸ 语音交互完善初步实现",
         "▸ 知识图谱智能搜索网关",
         "▸ 新增Text2Cypher模块",
         "  LLM动态生成Cypher查询",
         "",
         "▸ Cypher注入修复（安全优先）",
         "",
         "",
         "",
         "",
         "",
     ], icon="🕸️", title_color=PRIMARY)

card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "具体负责模块（请人工补充）", [
         "▸ Neo4j知识图谱数据构建",
         "",
         "▸ [请补充其他负责内容]",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
     ], icon="📝", title_color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 16 — 个人分工（林学成）
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "个人分工", "PART 04 · 林学成")

card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "CI/CD与基础设施", [
         "▸ A线改进——基础设施与实时通信全面升级",
         "▸ CI流水线修复",
         "  使用轻量依赖避免磁盘/导入崩溃",
         "▸ conftest.py延迟加载settings",
         "  避免CI缺.env时报错",
         "▸ 放宽CI依赖版本约束",
         "  移除不需要的包",
         "",
         "▸ 优化Agent智能层",
         "  LLM消息格式重构、多轮对话记忆、",
         "  工具调用系统",
     ], icon="🚀", title_color=PRIMARY)

card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "具体负责模块（请人工补充）", [
         "▸ [请补充其他负责内容]",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
         "",
     ], icon="📝", title_color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 17 — 项目总结：成果与亮点
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "项目总结", "PART 05 · 成果与亮点")

highlights = [
    ("🎯 系统架构", "完成从单体到多智能体架构的演进\n21个API端点 · 20个服务模块 · 5个专业Agent"),
    ("🧠 AI能力", "集成DeepSeek LLM + 知识图谱 + NER\nText2Cypher · 异常检测 · 多模态交互"),
    ("📋 功能覆盖", "全周期健康管理：评估·打卡·康复·预警\n患者端+医生端双角色完整覆盖"),
    ("⚡ 工程质量", "CI/CD流水线 · Docker容器化 · 单元测试\n结构化日志 · Prometheus监控"),
    ("🔒 安全保障", "JWT认证+Token黑名单 · Cypher注入防护\n连接池优化 · 统一异常处理"),
    ("📱 用户体验", "Vue 3响应式 · SSE流式 · WebSocket实时\nECharts图表 · 成就激励系统"),
]
for i, (title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    x = Inches(0.6) + Inches(col * 4.1)
    y = Inches(1.4) + Inches(row * 2.8)
    shp = rect(s, x, y, Inches(3.8), Inches(2.4), WHITE)
    rect(s, x, y, Inches(0.06), Inches(2.4), GREEN)
    txt(s, x + Inches(0.25), y + Inches(0.2), Inches(3.3), Inches(0.35),
        title, size=16, bold=True, color=GREEN)
    txt(s, x + Inches(0.25), y + Inches(0.65), Inches(3.3), Inches(1.6),
        desc, size=12, color=GRAY)


# ══════════════════════════════════════════════════════════
#  Slide 18 — 项目总结：未来规划
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "项目总结", "PART 05 · 未来规划")

plans = [
    ("🔄 持续优化", "模型微调优化 · 更多医学领域知识\n提升NER准确率 · 增强意图识别能力"),
    ("📱 移动端扩展", "微信小程序 · Flutter移动应用\n离线模式支持 · 推送通知优化"),
    ("🌐 数据分析", "患者群体健康趋势分析\n康复效果评估报告 · 医疗决策支持"),
    ("🤝 医疗集成", "对接医院HIS系统 · 电子病历互通\n远程会诊支持 · 处方审核"),
]
for i, (title, desc) in enumerate(plans):
    x = Inches(0.6) + Inches(i * 3.15)
    shp = rect(s, x, Inches(1.5), Inches(2.95), Inches(2.8), LIGHT_BLUE)
    shp.line.fill.background()
    txt(s, x + Inches(0.2), Inches(1.7), Inches(2.55), Inches(0.4),
        title, size=16, bold=True, color=PRIMARY)
    txt(s, x + Inches(0.2), Inches(2.2), Inches(2.55), Inches(1.8),
        desc, size=12, color=GRAY)

# 致谢
txt(s, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.8),
    "感谢聆听！", size=36, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.5),
    "术后全周期健康管理系统 — 让术后康复更智能、更高效", size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════
#  Slide 19 — 截图操作指南（供演示用，可删除）
# ══════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, WHITE)
title_bar(s, "附录：截图演示指南", "APPENDIX · 截图操作说明")

card(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5),
     "需要截图的页面（共6组）", [
         "① 患者端AI对话界面",
         "   → 展示多轮对话、NER识别效果",
         "",
         "② 知识图谱查询界面",
         "   → 展示Text2Cypher、实体关系图",
         "",
         "③ 医生端仪表盘",
         "   → 展示高风险患者、告警列表",
         "",
         "④ 康复计划界面",
         "   → 展示计划生成、指标追踪、成就",
         "",
         "⑤ 多模态交互",
         "   → 展示语音/图片/OCR功能",
         "",
         "⑥ 每日健康打卡",
         "   → 展示打卡界面与数据录入",
     ], icon="📸", title_color=PRIMARY)

card(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
     "截图操作步骤", [
         "1. 启动前后端服务",
         "   后端: python run.py",
         "   前端: npm run dev",
         "",
         "2. 浏览器打开 http://localhost:5173",
         "",
         "3. 分别以患者/医生账号登录",
         "",
         "4. 导航到对应功能页面",
         "",
         "5. 使用截图工具截取",
         "   推荐: Win+Shift+S 或 Snipaste",
         "",
         "6. 截图保存后替换PPT中的占位区域",
     ], icon="🔧", title_color=SECONDARY)


# ══════════════════════════════════════════════════════════
#  保存
# ══════════════════════════════════════════════════════════
output_path = r"E:\College\26春 软件工程创新实践\全周期健康管理系统汇报.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
